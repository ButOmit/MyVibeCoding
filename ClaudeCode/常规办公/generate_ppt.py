#!/usr/bin/env python3
"""Generate an academic PPT on All-Perovskite Tandem Solar Cells (~10 min talk)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── Color Palette (dark academic blue theme) ──
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
BG_SLIDE   = RGBColor(0x24, 0x24, 0x3D)   # slightly lighter
ACCENT     = RGBColor(0x00, 0xBF, 0xFF)   # cyan accent
ACCENT2    = RGBColor(0xFF, 0x6B, 0x35)   # orange accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)
GREEN      = RGBColor(0x00, 0xE6, 0x76)

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# ── Helpers ──
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="", font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multiline_box(slide, left, top, width, height, lines, font_size=16,
                      color=WHITE, bold_first=True, line_spacing=1.5, font_name="Microsoft YaHei"):
    """lines is a list of (text, is_bold) tuples."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bld if bold_first and i == 0 else bld
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return tf

def add_accent_bar(slide, left, top, width, height=0.06, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_page_number(slide, num, total=12):
    add_textbox(slide, 12.0, 7.0, 1.2, 0.4, f"{num} / {total}",
                font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, title, subtitle=None):
    add_accent_bar(slide, 0.8, 1.0, 2.0, 0.06)
    add_textbox(slide, 0.8, 1.15, 11.5, 0.7, title, font_size=30, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, 0.8, 1.75, 11.5, 0.5, subtitle, font_size=14, color=LIGHT_GRAY)

def add_bullet_points(slide, left, top, width, height, bullets, font_size=15, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(8)
        p.level = 0
    return tf

def add_card(slide, left, top, width, height, title, content_lines, title_color=ACCENT, card_color=None):
    """Add a card-like box with title and content."""
    if card_color is None:
        card_color = RGBColor(0x2A, 0x2A, 0x48)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = card_color
    shape.line.fill.background()
    # round corners
    shape.adjustments[0] = 0.05

    add_textbox(slide, left + 0.25, top + 0.15, width - 0.5, 0.4,
                title, font_size=16, color=title_color, bold=True)
    txBox = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.55),
                                     Inches(width - 0.5), Inches(height - 0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(4)
    return shape

def create_divider_slide(title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)
    # large centered title
    add_textbox(slide, 1.5, 2.5, 10.3, 1.5, title_text,
                font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_accent_bar(slide, 5.5, 3.8, 2.3, 0.06, color=ACCENT)
    if subtitle_text:
        add_textbox(slide, 1.5, 4.1, 10.3, 0.8, subtitle_text,
                    font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    return slide


# =====================================================
# SLIDE 1: TITLE
# =====================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide1, BG_DARK)

# Decorative top bar
add_accent_bar(slide1, 0, 0, 13.333, 0.08, ACCENT)
add_accent_bar(slide1, 0, 7.42, 13.333, 0.08, ACCENT)

# Title
add_textbox(slide1, 1.5, 1.8, 10.3, 1.2,
            "全钙钛矿叠层太阳能电池技术",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide1, 1.5, 2.9, 10.3, 0.8,
            "All-Perovskite Tandem Solar Cells: Progress, Challenges & Perspectives",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide1, 5.0, 3.7, 3.3, 0.06, color=ACCENT)

# Author / Affiliation placeholder
add_textbox(slide1, 1.5, 4.2, 10.3, 0.5,
            "报告人  |  单位",
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide1, 1.5, 4.7, 10.3, 0.5,
            "2026年5月",
            font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom decorative elements
for i, (x, c) in enumerate([(4.0, ACCENT), (6.0, ACCENT2), (8.0, GOLD)]):
    shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(6.0), Inches(1.0), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()


# =====================================================
# SLIDE 2: OUTLINE
# =====================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)
add_page_number(slide2, 2)

add_section_title(slide2, "报告提纲", "Outline")

outline_items = [
    ("01", "研究背景与意义", "光伏发展现状与钙钛矿材料的崛起"),
    ("02", "单结太阳能电池的效率极限", "Shockley-Queisser极限与叠层突破思路"),
    ("03", "全钙钛矿叠层电池结构与原理", "宽带隙/窄带隙子电池设计"),
    ("04", "关键技术挑战", "相稳定性、锡铅混合、互联层、大面积制备"),
    ("05", "最新研究进展与效率纪录", "2024-2026里程碑成果"),
    ("06", "稳定性与规模化前景", "商业化挑战与未来展望"),
]

for i, (num, title, desc) in enumerate(outline_items):
    y = 2.3 + i * 0.8
    add_textbox(slide2, 1.2, y, 0.7, 0.5, num, font_size=22, color=ACCENT, bold=True)
    add_textbox(slide2, 2.0, y, 3.5, 0.5, title, font_size=16, color=WHITE, bold=True)
    add_textbox(slide2, 5.5, y + 0.05, 6.5, 0.5, desc, font_size=12, color=LIGHT_GRAY)


# =====================================================
# SLIDE 3: RESEARCH BACKGROUND
# =====================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)
add_page_number(slide3, 3)

add_section_title(slide3, "研究背景：光伏技术的快速发展",
                  "Research Background: Rapid Development of Photovoltaics")

# Key data cards
add_card(slide3, 0.8, 2.3, 3.7, 2.2, "全球光伏装机量 (2025)",
         ["• 累计装机突破 2 TWp", "• 年新增装机 > 500 GW",
          "• 中国占全球新增装机 ~50%", "• 光伏已成最便宜电力来源之一"],
         title_color=ACCENT, card_color=RGBColor(0x2A, 0x2A, 0x48))

add_card(slide3, 4.8, 2.3, 3.7, 2.2, "晶硅电池的瓶颈",
         ["• 实验室效率接近 27% 极限", "• 单结理论极限 ~29.4% (S-Q)",
          "• 效率提升空间日益收窄", "• 亟需新一代高效率低成本技术"],
         title_color=ACCENT2, card_color=RGBColor(0x2A, 0x2A, 0x48))

add_card(slide3, 8.8, 2.3, 3.7, 2.2, "钙钛矿：变革性材料",
         ["• 2009年效率 3.8% → 2025年 >26%", "• 带隙可调 1.2-2.3 eV",
          "• 溶液法制备，成本低", "• 缺陷容忍度高，开路电压高"],
         title_color=GOLD, card_color=RGBColor(0x2A, 0x2A, 0x48))

# Bottom highlight
add_textbox(slide3, 0.8, 5.0, 11.5, 0.6,
            "▎核心思路：超越单结Shockley-Queisser极限 → 叠层电池 → 全钙钛矿叠层（成本最低的叠层路线）",
            font_size=14, color=ACCENT, bold=True)

# Efficiency evolution
add_textbox(slide3, 0.8, 5.8, 11.5, 1.0,
            "钙钛矿效率演进：  2009: 3.8% → 2016: 22.1% → 2020: 25.5% → 2023: 26.1% → 2025: 26.7% (认证)",
            font_size=13, color=LIGHT_GRAY)

# Bottom mini citation
add_textbox(slide3, 0.8, 6.8, 11.5, 0.4,
            "NREL Best Research-Cell Efficiency Chart, 2025; J. Phys. Chem. Lett. 2013, 4, 3623; Science 2024, 386, 1028",
            font_size=9, color=LIGHT_GRAY)


# =====================================================
# SLIDE 4: S-Q LIMIT & TANDEM MOTIVATION
# =====================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_DARK)
add_page_number(slide4, 4)

add_section_title(slide4, "从单结到叠层：突破 S-Q 极限",
                  "Beyond Single Junction: Breaking the Shockley-Queisser Limit")

# Left: S-Q limit explanation
add_card(slide4, 0.8, 2.3, 5.5, 4.2,
         "Shockley-Queisser 极限 (~1961)",
         ["• 单结太阳能电池的理论效率上限：~33.7%",
          "  (最佳带隙 ~1.34 eV)",
          "",
          "两大损失机制：",
          "  ① 亚带隙光子透射损失 (Below-Eg loss)",
          "     → 低能光子不被吸收，能量浪费",
          "  ② 热化损失 (Thermalization loss)",
          "     → 高能光子多余能量以热耗散",
          "",
          "单结 Si 电池极限：~29.4%",
          "当前实验室纪录：~27.3% (LONGi, 2024)",
          "晶硅电池已逼近其理论天花板"],
         title_color=ACCENT2)

# Right: Tandem principle
add_card(slide4, 6.8, 2.3, 5.7, 4.2,
         "叠层策略：光谱分光利用",
         ["核心思想：",
          "  宽带隙顶电池吸收高能光子 (紫外-可见)",
          "  窄带隙底电池吸收低能光子 (近红外)",
          "  → 减小热化损失 + 利用更多光谱",
          "",
          "理论效率上限：",
          "  • 双结叠层：~45%（极限）",
          "  • 三结叠层：~50%（极限）",
          "",
          "叠层电池分类：",
          "  • III-V 族 (GaAs/GaInP)：效率最高，成本极高",
          "  • 钙钛矿/晶硅叠层：效率 > 34% (2025)",
          "  • 全钙钛矿叠层：全部溶液法制备，成本最低 !"],
         title_color=ACCENT)

# Key takeaway
add_textbox(slide4, 0.8, 6.7, 11.5, 0.5,
            "▎全钙钛矿叠层 = 最低成本路线 + 高效率潜力 + 柔性/轻质应用前景",
            font_size=15, color=GOLD, bold=True)


# =====================================================
# SLIDE 5: ALL-PEROVSKITE TANDEM STRUCTURE
# =====================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_DARK)
add_page_number(slide5, 5)

add_section_title(slide5, "全钙钛矿叠层电池：器件结构",
                  "Device Architecture of All-Perovskite Tandem Solar Cells")

# Device stack - vertical text representation
layers = [
    ("Glass / 柔性基底", LIGHT_GRAY, None),
    ("ITO 透明电极", WHITE, None),
    ("空穴传输层 (HTL)", RGBColor(0x66, 0xCC, 0xFF), "NiOx, PTAA, MeO-2PACz"),
    ("宽带隙钙钛矿顶电池", ACCENT, "~1.75-1.85 eV  (吸收 300-700 nm)"),
    ("电子传输层 (ETL)", RGBColor(0x66, 0xCC, 0xFF), "C₆₀, PCBM, SnO₂"),
    ("隧穿复合层 / 互联层", RGBColor(0xFF, 0x99, 0x66), "ITO / AZO / 超薄金属"),
    ("空穴传输层 (HTL)", RGBColor(0x66, 0xCC, 0xFF), "PEDOT:PSS, NiOx"),
    ("窄带隙钙钛矿底电池", ACCENT2, "~1.20-1.30 eV  (吸收 700-1100 nm)"),
    ("电子传输层 (ETL)", RGBColor(0x66, 0xCC, 0xFF), "C₆₀ / BCP"),
    ("金属电极 (Ag/Cu/Au)", WHITE, None),
]

# Left side: stack diagram
y_start = 2.0
layer_h = 0.44
for i, (name, color, detail) in enumerate(layers):
    y = y_start + i * layer_h
    shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(1.0), Inches(y),
                                    Inches(5.0), Inches(layer_h - 0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    alpha = 0.9 if i in [3, 7] else 0.6
    # text label on bar
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(12)
    p.font.color.rgb = BG_DARK if color != BG_DARK else WHITE
    p.font.bold = (i in [3, 7])
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(2)

    # detail text
    if detail:
        add_textbox(slide5, 6.3, y - 0.02, 5.5, 0.4, detail, font_size=11, color=LIGHT_GRAY)

# Labels
add_textbox(slide5, 1.0, 1.6, 5.0, 0.4, "☀ 太阳光入射方向 ↓", font_size=11, color=GOLD, alignment=PP_ALIGN.CENTER)
add_textbox(slide5, 9.0, 6.2, 3.5, 0.4, "• 共 10+ 层薄膜堆叠", font_size=10, color=LIGHT_GRAY)
add_textbox(slide5, 9.0, 6.5, 3.5, 0.4, "• 总厚度 < 3 μm", font_size=10, color=LIGHT_GRAY)
add_textbox(slide5, 9.0, 6.8, 3.5, 0.4, "• 溶液法 + 蒸镀法结合", font_size=10, color=LIGHT_GRAY)

# Right side: key parameters
add_card(slide5, 6.3, 7.0, 6.2, 0,  # placeholder - we handle text manually
         "", [])
# Actually let's add parameter boxes
params = [
    ("顶电池 Eg", "1.75-1.85 eV", "吸收 300-700 nm", ACCENT),
    ("底电池 Eg", "1.20-1.30 eV", "吸收 700-1100 nm", ACCENT2),
    ("理想 Jsc", "~15-16 mA/cm²", "电流匹配条件", GOLD),
    ("理论 Voc", "> 2.1 V", "双结叠加优势", GREEN),
]
for i, (label, value, desc, color) in enumerate(params):
    y = 2.0 + i * 1.2
    add_textbox(slide5, 11.8, y, 1.2, 0.3, label, font_size=10, color=color, bold=True)
    add_textbox(slide5, 11.8, y + 0.25, 1.2, 0.4, value, font_size=16, color=WHITE, bold=True)
    add_textbox(slide5, 11.8, y + 0.62, 1.2, 0.3, desc, font_size=9, color=LIGHT_GRAY)


# =====================================================
# SLIDE 6: WIDE-BANDGAP PEROVSKITE
# =====================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, BG_DARK)
add_page_number(slide6, 6)

add_section_title(slide6, "宽带隙钙钛矿：高 Br 含量引发的挑战",
                  "Wide-Bandgap Perovskite: Phase Segregation & Voltage Loss")

add_card(slide6, 0.8, 2.2, 5.8, 2.0,
         "带隙调节策略",
         ["• ABX₃ 结构中 X 位卤素混合调控 Eg",
          "  MAPb(I₁₋ₓBrₓ)₃:  x 从 0→1, Eg 从 1.55→2.30 eV",
          "• 叠层顶电池目标：Eg ~1.75-1.85 eV → x~0.3-0.4",
          "• 问题：高 Br 含量 → 光致卤素相分离 !"],
         title_color=ACCENT2)

add_card(slide6, 0.8, 4.5, 5.8, 2.5,
         "光致卤素相分离 (Hoke Effect)",
         ["• 光照下 I-rich 和 Br-rich 相分离",
          "  MAPb(I₀.₆Br₀.₄)₃ → I-rich (Eg↓) + Br-rich (Eg↑)",
          "• 后果：",
          "  ① 低 Eg 相成为载流子陷阱/复合中心",
          "  ② Voc 持续下降，性能退化",
          "  ③ 带隙分布不均，光谱响应畸变",
          "• 机理：极化子诱导 / 晶格应变 / 卤素空位迁移"],
         title_color=ACCENT2)

add_card(slide6, 7.0, 2.2, 5.5, 4.8,
         "缓解策略",
         ["1. A 位阳离子工程",
          "   • Cs⁺/FA⁺/MA⁺ 混合 → 抑制离子迁移",
          "   • 如 Cs₀.₂FA₀.₈Pb(I₀.₆Br₀.₄)₃",
          "",
          "2. 添加剂工程",
          "   • K⁺, Rb⁺ 钝化晶界缺陷",
          "   • PEAI, BAI 等大阳离子 → 2D/3D 异质结",
          "   • MACl 辅助结晶",
          "",
          "3. 三维/低维异质结构",
          "   • 表面 2D 层 → 抑制卤素迁移通道",
          "   • 形成能量势垒，稳定体相组成",
          "",
          "4. 晶粒尺寸与应变工程",
          "   • 大晶粒 → 减少晶界迁移路径",
          "   • 组分梯度钝化"],
         title_color=ACCENT)


# =====================================================
# SLIDE 7: NARROW-BANDGAP Sn-Pb PEROVSKITE
# =====================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, BG_DARK)
add_page_number(slide7, 7)

add_section_title(slide7, "窄带隙 Sn-Pb 钙钛矿：低带隙的关键",
                  "Narrow-Bandgap Perovskite: The Sn-Pb Challenge")

add_card(slide7, 0.8, 2.2, 5.8, 2.2,
         "Sn-Pb 混合策略",
         ["• 纯 Pb 钙钛矿最低 Eg ~1.48 eV (FAPbI₃)",
          "  → 对叠层底电池来说带隙偏高",
          "• 引入 Sn 替代 Pb → 带隙进一步降低",
          "  MASnₓPb₁₋ₓI₃: x↑ → Eg↓ (最低 ~1.20 eV)",
          "• 最佳底电池 Eg: 1.20-1.25 eV",
          "• 理想 Sn 含量：25-50% Pb 替代"],
         title_color=ACCENT2)

add_card(slide7, 0.8, 4.7, 5.8, 2.2,
         "核心挑战：Sn²⁺ 氧化",
         ["• Sn²⁺ → Sn⁴⁺ 极易氧化 (空气中数分钟)",
          "  → 形成 Sn 空位 (p 型掺杂严重)",
          "  → 载流子寿命短，Voc 损失大",
          "",
          "• 解决策略：",
          "  ① SnF₂ / SnCl₂ 添加剂 (提供过量 Sn²⁺)",
          "  ② 还原性添加剂 (Sn⁰ 粉, 金属 Sn)",
          "  ③ 抗氧化剂 (抗坏血酸, HPA, etc.)",
          "  ④ 惰性气氛下制备 (N₂ glovebox)"],
         title_color=ACCENT2)

add_card(slide7, 7.0, 2.2, 5.5, 4.7,
         "最新突破",
         ["1. 组分优化",
          "   • FA₀.₇MA₀.₃Pb₀.₅Sn₀.₅I₃ 基准组分",
          "   • Cs⁺ / Gua⁺ 稳定钙钛矿相",
          "",
          "2. 添加剂协同",
          "   • SnF₂ + 金属 Sn 粉 + NH₄SCN",
          "   • 大幅提升载流子寿命 (>1 μs)",
          "",
          "3. 2D/3D 异质结钝化",
          "   • PEAI / BAI 表面处理",
          "   • 减少表面 Sn 空位缺陷",
          "",
          "4. 性能里程碑",
          "   • 单结 Sn-Pb PSC: > 23%",
          "   • 载流子扩散长度 > 3 μm",
          "   • Voc deficit < 0.40 V"],
         title_color=ACCENT)


# =====================================================
# SLIDE 8: INTERCONNECTION LAYER
# =====================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, BG_DARK)
add_page_number(slide8, 8)

add_section_title(slide8, "互联层 (ICL) 与隧穿复合结",
                  "Interconnection Layer: The Critical Junction")

add_card(slide8, 0.8, 2.2, 5.5, 2.3,
         "互联层的核心功能",
         ["• 串联子电池的关键界面",
          "• 要求：",
          "  ① 低电阻欧姆接触 → 高效载流子复合",
          "  ② 高透光率 → 不影响底电池光吸收",
          "  ③ 化学惰性 → 不破坏下方钙钛矿层",
          "  ④ 低横向电导 → 避免分流",
          "• 本质：电子/空穴的隧穿复合中心"],
         title_color=ACCENT)

add_card(slide8, 6.8, 2.2, 5.7, 2.3,
         "常见 ICL 体系",
         ["• 超薄金属层：Au (~1 nm) / Ag",
          "  优点：简单  |  缺点：寄生吸收",
          "• 透明导电氧化物 (TCO)：",
          "  ITO (溅射), AZO (ALD), IZO",
          "  → 目前主流方案，光学/电学兼顾",
          "• 复合层 + TCO 叠层：",
          "  C₆₀ / SnO₂ / ITO 或 PEI / ITO",
          "• 新兴方案：",
          "  MoOₓ / ITO, 超薄 Ag / AZO"],
         title_color=ACCENT2)

add_card(slide8, 0.8, 4.8, 11.7, 2.2,
         "ICL 工程的关键进展",
         ["• 溅射损伤问题 → ALD 低温生长 ITO/AZO，保护下方有机层",
          "• 光学管理 → 调控 ICL 厚度优化干涉效应，增强底电池光吸收（光学 spacer 效应）",
          "• 复合效率 → 通过掺杂调控功函数，确保 ETL(顶)/HTL(底) 能带对齐",
          "• 大面积均匀性 → 原子层沉积 ALD 保证 pinhole-free，实现 cm² 级均匀隧穿"],
         title_color=GOLD)


# =====================================================
# SLIDE 9: EFFICIENCY PROGRESS
# =====================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9, BG_DARK)
add_page_number(slide9, 9)

add_section_title(slide9, "全钙钛矿叠层效率纪录演进",
                  "Efficiency Progress of All-Perovskite Tandems")

# Timeline entries
timeline = [
    ("2016", "14.1%", "Stanford / MIT\n首个全钙钛矿四端叠层"),
    ("2019", "23.1%", "NREL\nSn-Pb 底电池突破"),
    ("2020", "24.8%", "南京大学\nFA-Cs 宽带隙+Sn-Pb"),
    ("2021", "26.4%", "KAUST\n组分+添加剂协同优化"),
    ("2022", "28.0%", "南京大学\n大面积认证效率突破"),
    ("2023", "29.1%", "KAUST / UNIST\n>29% 认证效率里程碑"),
    ("2024", "30.2%", "南京大学\n首超30% 全钙钛矿叠层"),
    ("2025", ">30.5%", "多课题组\n大面积+柔性方向并进"),
]

# Draw timeline
y_base = 2.2
add_accent_bar(slide9, 0.8, y_base + 0.15, 11.5, 0.04, ACCENT)

for i, (year, eff, desc) in enumerate(timeline):
    x = 0.8 + i * 1.45
    # Dot on timeline
    dot = slide9.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x + 0.3), Inches(y_base + 0.05),
                                  Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT if i < 7 else GOLD
    dot.line.fill.background()

    # Year
    add_textbox(slide9, x, y_base - 0.55, 1.0, 0.35, year, font_size=13, color=ACCENT, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Efficiency
    add_textbox(slide9, x, y_base + 0.35, 1.0, 0.45, eff, font_size=20, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide9, x - 0.2, y_base + 0.75, 1.4, 0.9, desc, font_size=8, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)

# Arrow showing progress
add_textbox(slide9, 0.8, 1.6, 11.5, 0.5,
            "效率从 14.1% → 30.5%+ ，仅用 10 年！增速超越晶硅与钙钛矿/晶硅叠层",
            font_size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# Comparison note
add_card(slide9, 0.8, 5.0, 5.5, 2.0,
         "效率对比 (2025)",
         ["• 晶硅单结:      ~27.3% (LONGi)",
          "• 钙钛矿单结:    ~26.7% (认证)",
          "• 钙钛矿/晶硅叠层: ~34.6% (LONGi)",
          "• 全钙钛矿叠层:  ~30.5%+",
          "→ 全钙钛矿在无硅基底情况下已经突破30%"],
         title_color=ACCENT)

add_card(slide9, 6.8, 5.0, 5.7, 2.0,
         "关键特征",
         ["• 全溶液法制备，成本远低于晶硅/III-V族叠层",
          "• 可制备在柔性基底上 → 轻质、柔性应用",
          "• 理论效率上限 ~45%，仍有巨大提升空间",
          "• 大面积 (cm²级) 效率正在追赶小面积纪录",
          "• 中国在领域内处于世界领先地位 (南大、北大等)"],
         title_color=GOLD)


# =====================================================
# SLIDE 10: STABILITY & SCALABILITY
# =====================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10, BG_DARK)
add_page_number(slide10, 10)

add_section_title(slide10, "稳定性与规模化：从实验室到产业化",
                  "Stability & Scalability: From Lab to Fab")

add_card(slide10, 0.8, 2.2, 3.7, 2.4,
         "稳定性挑战",
         ["• 湿度/氧气 → 钙钛矿分解",
          "• 热应力 → 相变、组分挥发",
          "  (MA⁺, I⁻ 热稳定性差)",
          "• 光照 → 相分离 (WBG)",
          "  Sn²⁺氧化 (NBG)",
          "• 电场 → 离子迁移",
          "• 电极腐蚀 → ITO-In, Ag-I 反应",
          "• 封装是必须的"],
         title_color=ACCENT2)

add_card(slide10, 4.8, 2.2, 3.7, 2.4,
         "稳定性提升策略",
         ["• 全无机 / Cs-FA 体系替代 MA",
          "• 2D/3D 异质结 → 防潮层",
          "• 离子液体添加剂 (IL)",
          "• 碳基电极替代金属电极",
          "• ALD 封装 (Al₂O₃, SiNₓ)",
          "• 自修复聚合物封装",
          "• IEC 61215 标准测试",
          "  (湿热/热循环/UV 老化)"],
         title_color=ACCENT)

add_card(slide10, 8.8, 2.2, 3.7, 2.4,
         "大面积制备方法",
         ["• 旋涂法 → 刮涂 (Blade coating)",
          "• 狭缝涂布 (Slot-die coating)",
          "• 喷涂 (Spray coating)",
          "• 气相沉积 (蒸镀 + CVD)",
          "• 卷对卷 (R2R) 柔性制备",
          "• 模块化 (Module) 设计",
          "  → 消除死区面积损失",
          "• 激光划线 P1-P2-P3 工艺"],
         title_color=GOLD)

# Bottom highlight
add_card(slide10, 0.8, 4.9, 11.7, 2.2,
         "产业化现状与展望",
         ["• 领先企业已开始中试线建设 → 协鑫、极电光能、Oxford PV (钙钛矿/硅)、Swift Solar (全钙钛矿)",
          "• 2024-2025 年全球多家公司公布叠层组件量产计划，目标效率 > 25% 组件级",
          "• 全钙钛矿叠层独特优势 → 柔性轻质、可溶液法一体化制备，是唯一有望实现 < $0.10/W 成本的叠层路线",
          "• 关键产业化瓶颈：① 大面积均匀性  ② 长期户外稳定性 (>25 年)  ③ 铅毒性与回收  ④ 标准化测试认证"],
         title_color=GREEN)


# =====================================================
# SLIDE 11: KEY CHALLENGES SUMMARY
# =====================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide11, BG_DARK)
add_page_number(slide11, 11)

add_section_title(slide11, "核心挑战与应对策略总结",
                  "Summary of Key Challenges and Strategies")

challenges = [
    ("宽带隙钙钛矿\n光致相分离",
     "• 组分工程 (Cs/FA混合)\n• 2D/3D异质结\n• 添加剂钝化\n• 晶格应变调控",
     ACCENT),
    ("Sn-Pb钙钛矿\nSn²⁺氧化",
     "• 还原性添加剂\n• 惰性气氛制备\n• 表面钝化层\n• 全无机Cs基探索",
     ACCENT2),
    ("互联层\n光学/电学损耗",
     "• ALD-TCO 低温生长\n• 光学spacer设计\n• 能带工程调控\n• 超薄复合层",
     ACCENT),
    ("大面积\n均匀性",
     "• 刮涂/狭缝涂布\n• 气相法辅助\n• 激光划线优化\n• 在线监测反馈",
     ACCENT2),
]

for i, (challenge, solution, color) in enumerate(challenges):
    x = 0.8 + i * 3.1
    add_card(slide11, x, 2.2, 2.8, 2.0, "", [], title_color=color,
             card_color=RGBColor(0x2A, 0x2A, 0x48))
    add_textbox(slide11, x + 0.2, 2.35, 2.4, 1.0, challenge, font_size=15, color=color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide11, x + 0.2, 3.4, 2.4, 1.8, solution, font_size=11, color=LIGHT_GRAY)

# Efficiency loss breakdown
add_card(slide11, 0.8, 4.7, 11.7, 2.0,
         "效率损失通道分析与潜力",
         ["• Voc loss: 目前 Voc ~2.0-2.1 V (理论最大值 ~2.5 V) → 主要来自非辐射复合，改进空间 ~400 mV",
          "• FF loss: 目前 FF ~78-82% → 主要来自串联电阻与界面复合，改进空间 ~5-8 abs%",
          "• Jsc loss: 目前 Jsc ~15-16 mA/cm² → 反射/寄生吸收/电流不匹配，改进空间 ~2-3 mA/cm²",
          "• 综合预测：如解决上述 loss → 全钙钛矿叠层可达 ~35-38%，接近理论极限"],
         title_color=GOLD)


# =====================================================
# SLIDE 12: SUMMARY & OUTLOOK
# =====================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide12, BG_DARK)
add_page_number(slide12, 12)

add_section_title(slide12, "总结与展望", "Summary & Outlook")

# Summary
add_card(slide12, 0.8, 2.2, 5.8, 2.8,
         "总结",
         ["1. 全钙钛矿叠层电池是突破单结S-Q极限",
          "   最具成本优势的技术路线",
          "",
          "2. 效率已突破 30%，10 年提升 16 个百分点",
          "   → 增长速度为所有光伏技术中最快",
          "",
          "3. 核心挑战：宽带隙相分离、Sn-Pb氧化、",
          "   互联层工程、大面积制备",
          "",
          "4. 关键策略：组分/添加剂/界面/维度工程",
          "   多手段协同是提升性能的核心"],
         title_color=ACCENT)

add_card(slide12, 7.0, 2.2, 5.5, 2.8,
         "展望",
         ["1. 效率目标：2028 年前达到 35%+",
          "",
          "2. 稳定性：通过 IEC 61215 全套测试",
          "   → 户外寿命 > 25 年",
          "",
          "3. 大面积：从 cm² 到 m² 级组件",
          "   → 组件效率 > 25%",
          "",
          "4. 柔性全钙钛矿叠层 → 轻质光伏新应用",
          "   (BIPV, 可穿戴, 无人机, 航天)",
          "",
          "5. 全铅-free Sn 基叠层 → 环保终极方案"],
         title_color=GOLD)

# Bottom quote
add_textbox(slide12, 1.5, 5.5, 10.3, 1.0,
            '"全钙钛矿叠层太阳能电池代表了光伏技术的下一个前沿——\n'
            '它可能成为最廉价的超高效率太阳能电池技术路径。"',
            font_size=16, color=ACCENT2, bold=False, alignment=PP_ALIGN.CENTER)

add_textbox(slide12, 1.5, 6.5, 10.3, 0.5,
            "谢谢！欢迎提问  |  Thank you! Questions & Discussion",
            font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Decorative bottom bar
add_accent_bar(slide12, 0, 7.42, 13.333, 0.08, ACCENT)


# =====================================================
# SAVE
# =====================================================
output_path = r"E:\上学活动文件\碳中和全钙钛矿叠层pre\全钙钛矿叠层太阳能电池_学术报告.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
